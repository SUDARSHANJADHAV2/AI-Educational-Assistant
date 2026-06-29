import os
import json
import csv
import io
import requests
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import pandas as pd
from youtube_transcript_api import YouTubeTranscriptApi
import chromadb
from chromadb.utils import embedding_functions
from openai import AsyncOpenAI
from pypdf import PdfReader

# Initialize ChromaDB
chroma_client = chromadb.PersistentClient(path="./chroma_data")
emb_fn = embedding_functions.DefaultEmbeddingFunction()
collection = chroma_client.get_or_create_collection(name="studysnap_collection", embedding_function=emb_fn)

# Initialize LM Studio client
lm_studio_client = AsyncOpenAI(base_url="http://localhost:1234/v1", api_key="lm-studio")
MODEL_NAME = "my_custom_7B_model_gguf/Anish/StudySnap AI"

def extract_youtube_transcript(url: str) -> str:
    try:
        video_id = ""
        if "v=" in url:
            video_id = url.split("v=")[1].split("&")[0]
        elif "youtu.be/" in url:
            video_id = url.split("youtu.be/")[1].split("?")[0]
        elif "embed/" in url:
            video_id = url.split("embed/")[1].split("?")[0]
            
        if not video_id: 
            raise ValueError("Could not find video ID in URL.")
        
        ytt_api = YouTubeTranscriptApi()
        
        try:
            # Try getting standard English or auto-generated
            transcript_list = ytt_api.list(video_id)
            transcript = transcript_list.find_transcript(['en', 'en-US', 'en-GB', 'hi'])
            data = transcript.fetch()
        except:
            try:
                # Fallback to any auto-generated transcript if manual English isn't found
                transcript_list = ytt_api.list(video_id)
                transcript = list(transcript_list)[0] # Just get the first available one
                data = transcript.fetch()
            except Exception as e:
                # If all else fails, use the basic method
                data = ytt_api.fetch(video_id)
                
        return " ".join([t.text if hasattr(t, 'text') else t.get('text', '') for t in data])
    except Exception as e:
        print(f"YouTube Extraction Error: {e}")
        raise ValueError(f"Failed to process YouTube transcript: {str(e)}. The video might not have captions enabled.")

def extract_webpage(url: str) -> str:
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers)
        soup = BeautifulSoup(response.text, 'html.parser')
        # Remove scripts and styles
        for script in soup(["script", "style"]):
            script.extract()
        return soup.get_text(separator="\n", strip=True)
    except Exception as e:
        print(f"Webpage Extraction Error: {e}")
        return ""

import whisper

def extract_audio(file_path: str) -> str:
    """Extract text from audio using local Whisper instance."""
    try:
        # Load the base model ('tiny', 'base', 'small', 'medium', 'large')
        # 'base' is a good tradeoff between speed and accuracy for local transcription
        model = whisper.load_model("base")
        
        # Transcribe the audio file
        # Whisper automatically handles chunking for long audio files
        result = model.transcribe(file_path)
        
        return result.get("text", "")
    except Exception as e:
        print(f"Whisper Audio Extraction Error: {e}")
        return "Audio transcription failed. Please ensure ffmpeg is installed."

def extract_text(file_path: str, filename: str) -> str:
    """Extract text from all supported file types."""
    text = ""
    ext = filename.split('.')[-1].lower() if '.' in filename else ""
    
    try:
        if ext == 'pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                if page_text := page.extract_text(): text += page_text + "\n"
        elif ext in ['txt', 'md', 'html', 'xml']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                if ext in ['html', 'xml']:
                    soup = BeautifulSoup(content, 'html.parser')
                    text = soup.get_text(separator="\n", strip=True)
                else:
                    text = content
        elif ext == 'json':
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                text = json.dumps(data, indent=2)
        elif ext == 'csv':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    text += " ".join(row) + "\n"
        elif ext == 'xlsx':
            df = pd.read_excel(file_path)
            text = df.to_string()
        elif ext == 'docx':
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == 'pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ['wav', 'mp3', 'm4a', 'flac']:
            text = extract_audio(file_path)
    except Exception as e:
        print(f"Error extracting {filename}: {e}")
        
    return text

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200):
    chunks = []
    start = 0
    text_length = len(text)
    while start < text_length:
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
    return chunks

def ingest_document(file_path: str, filename: str, file_id: str):
    text = ""
    # Check if it's a URL or YouTube link based on filename format or type passed from main
    if filename.startswith("http"):
        if "youtube.com" in filename or "youtu.be" in filename:
            text = extract_youtube_transcript(filename)
            filename = f"YouTube: {filename}"
        else:
            text = extract_webpage(filename)
            filename = f"Web: {filename}"
    else:
        text = extract_text(file_path, filename)
        
    if not text.strip():
        raise ValueError("Could not extract any text from the document.")

    chunks = chunk_text(text)
    ids = [f"{file_id}_chunk_{i}" for i in range(len(chunks))]
    metadatas = [{"source": filename, "file_id": file_id} for _ in chunks]
    
    collection.add(documents=chunks, metadatas=metadatas, ids=ids)
    return {"id": file_id, "filename": filename, "type": filename.split('.')[-1] if '.' in filename and not filename.startswith("http") else "url"}

def get_all_sources():
    results = collection.get(include=["metadatas"])
    metadatas = results.get("metadatas", [])
    unique_sources = {}
    for meta in metadatas:
        if meta["file_id"] not in unique_sources:
            unique_sources[meta["file_id"]] = {
                "id": meta["file_id"],
                "filename": meta["source"],
                "type": meta["source"].split('.')[-1] if '.' in meta["source"] and not meta["source"].startswith("http") else "url"
            }
    return list(unique_sources.values())

def delete_source(file_id: str):
    """Delete a source and all its chunks from ChromaDB."""
    collection.delete(where={"file_id": file_id})

async def chat_with_context(query: str, history: list = None, selected_source_ids: list = None):
    if history is None: history = []
    
    where_clause = None
    if selected_source_ids is not None and len(selected_source_ids) > 0:
        if len(selected_source_ids) == 1:
            where_clause = {"file_id": selected_source_ids[0]}
        else:
            where_clause = {"file_id": {"$in": selected_source_ids}}
            
    # If a filter is provided, pass it to where
    if where_clause:
        results = collection.query(query_texts=[query], n_results=4, where=where_clause)
    else:
        results = collection.query(query_texts=[query], n_results=4)
        
    documents = results.get("documents", [[]])[0]
    metadatas = results.get("metadatas", [[]])[0]
    
    context = ""
    sources_used = []
    for doc, meta in zip(documents, metadatas):
        context += f"--- Source: {meta['source']} ---\n{doc}\n\n"
        if meta['source'] not in sources_used:
            sources_used.append(meta['source'])
            
    system_prompt = (
        "You are StudySnap AI, a helpful and knowledgeable assistant. "
        "Use the provided context to answer the user's question. "
        "If the answer is not in the context, say that you don't know based on the provided documents.\n\n"
        "Context:\n" + context
    )
    
    messages = [{"role": "system", "content": system_prompt}]
    for msg in history:
        messages.append({"role": msg.get("role", "user"), "content": msg.get("content", "")})
    messages.append({"role": "user", "content": query})
    
    response = await lm_studio_client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        temperature=0.7,
        max_tokens=1024,
    )
    
    answer = response.choices[0].message.content
    return answer, sources_used
